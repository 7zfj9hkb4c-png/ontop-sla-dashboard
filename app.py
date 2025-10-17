# app.py — OnTop SLA Dashboard (Local Insights, v10)

import streamlit as st
import pandas as pd
import numpy as np
import io, os, re
import matplotlib.pyplot as plt
from pathlib import Path
from matplotlib.colors import ListedColormap, BoundaryNorm

st.set_page_config(page_title="OnTop SLA Dashboard (v10)", layout="wide")
SLA_TARGET_DEFAULT = 95.0

# ===================== HELPERS =====================
def enforce_int_week(df, col="week"):
    if df is None or df.empty or col not in df.columns:
        return df
    out = df.copy()
    out[col] = pd.to_numeric(out[col], errors="coerce").astype("Int64")
    out = out.dropna(subset=[col])
    out[col] = out[col].astype(int)
    return out

def scale_pct_if_needed(series):
    if series is None or len(series) == 0:
        return series
    s = pd.to_numeric(series, errors="coerce")
    try:
        if s.dropna().median() <= 1.5:
            return s * 100.0
        return s
    except Exception:
        return s

def load_clean_excel(file_bytes_or_path):
    def _read(x, sheet):
        try:
            return pd.read_excel(x, sheet_name=sheet)
        except Exception:
            return pd.DataFrame()
    x = io.BytesIO(file_bytes_or_path) if isinstance(file_bytes_or_path, (bytes, bytearray)) else file_bytes_or_path
    data = {
        "OnTop_SLA_Global": _read(x, "OnTop_SLA_Global"),
        "OnTop_FRT_Global": _read(x, "OnTop_FRT_Global"),
        "SLA_By_Team": _read(x, "SLA_By_Team"),
        "FRT_By_Team": _read(x, "FRT_By_Team"),
        "Tickets_Created": _read(x, "Tickets_Created"),
        "Tickets_Solved": _read(x, "Tickets_Solved"),
        "Tickets_Breached": _read(x, "Tickets_Breached"),
        "Tickets_Unbreached": _read(x, "Tickets_Unbreached"),
        "Tickets_Inflow": _read(x, "Tickets_Inflow"),
        "Tickets_Breach_Split": _read(x, "Tickets_Breach_Split"),
        "dim_team": _read(x, "dim_team"),
        "dim_week": _read(x, "dim_week"),
        "mrt_global": _read(x, "mrt_global"),
        "mrt_by_team": _read(x, "mrt_by_team"),
        "team_compliance": _read(x, "team_compliance"),
        "Summary": _read(x, "Summary"),
        "Category_By_Week": _read(x, "Category_By_Week"),
    }
    # si la hoja de categorías no existe, intentamos detectar otra
    if data["Category_By_Week"].empty:
        try:
            xl = pd.ExcelFile(x) if isinstance(x, io.BytesIO) else pd.ExcelFile(file_bytes_or_path)
            for name in xl.sheet_names:
                if any(re.search(pat, name, re.IGNORECASE) for pat in [r"category.*(week|wic)", r"(contact).*category"]):
                    data["Category_By_Week"] = pd.read_excel(xl, sheet_name=name)
                    break
        except Exception:
            pass
    return data

def normalize_category_df(df):
    """Normaliza a: category, year, week, tickets."""
    if df is None or df.empty:
        return pd.DataFrame(columns=["category","year","week","tickets"])
    c = df.copy()
    c.columns = [str(x).strip().lower() for x in c.columns]
    col_cat = next((x for x in c.columns if x in ["category","contact type","contact_type","type"]), None)
    col_year = next((x for x in c.columns if x in ["year","anio","año"]), None)
    col_week = next((x for x in c.columns if x in ["week","wic","iso_week"]), None)
    col_tk = next((x for x in c.columns if x in ["tickets","count","qty","volume"]), None)

    # matriz (categorías filas, semanas columnas)
    if col_cat and (col_week is None or col_tk is None):
        value_cols = [col for col in c.columns if re.fullmatch(r"\d{1,2}", str(col))]
        if value_cols:
            long = c.melt(id_vars=[col_cat], value_vars=value_cols, var_name="week", value_name="tickets")
            long["category"] = long[col_cat]
            long["tickets"] = pd.to_numeric(long["tickets"], errors="coerce")
            long["week"] = pd.to_numeric(long["week"], errors="coerce").astype("Int64")
            year_val = pd.to_numeric(c.get("year"), errors="coerce").dropna().unique()
            year = int(year_val[0]) if len(year_val)>0 else pd.Timestamp.today().year
            long["year"] = int(year)
            out = long[["category","year","week","tickets"]].dropna(subset=["week","tickets"])
            out["week"] = out["week"].astype(int)
            return out

    if not (col_cat and col_year and col_week and col_tk):
        return pd.DataFrame(columns=["category","year","week","tickets"])

    out = c.rename(columns={col_cat:"category", col_year:"year", col_week:"week", col_tk:"tickets"})
    out["tickets"] = pd.to_numeric(out["tickets"], errors="coerce")
    out["week"] = pd.to_numeric(out["week"], errors="coerce").astype("Int64")
    out = out.dropna(subset=["week","tickets"])
    out["week"] = out["week"].astype(int)
    out["year"] = pd.to_numeric(out["year"], errors="coerce").fillna(pd.Timestamp.today().year).astype(int)
    return out[["category","year","week","tickets"]]

def infer_target(summary_df, default_target=SLA_TARGET_DEFAULT):
    if isinstance(summary_df, pd.DataFrame) and not summary_df.empty:
        s = summary_df.copy()
        s.columns = [str(c).strip().lower() for c in s.columns]
        if {"kpi","value"}.issubset(s.columns):
            row = s.loc[s["kpi"].astype(str).str.lower().eq("target_sla_pct"), "value"]
            if not row.empty:
                try:
                    return float(str(row.iloc[0]).replace("%","").replace(",","."))  # tolerante
                except:
                    pass
    return float(default_target)

def wow_delta(series, current_week):
    if series is None or series.empty or current_week not in series.index:
        return (np.nan, np.nan, np.nan)
    curr = series.get(current_week, np.nan)
    prev = series.get(current_week-1, np.nan)
    delta = (curr - prev) if pd.notna(curr) and pd.notna(prev) else np.nan
    return (curr, prev, delta)

def annotate_bars(ax, xs, ys):
    for x, y in zip(xs, ys):
        ax.text(x, (y/2 if y>0 else 0), f"{int(y):,}", ha="center", va="center")

def annotate_stacked_bars(ax, xs, bottom, heights):
    for x, b, h in zip(xs, bottom, heights):
        if h > 0:
            ax.text(x, b + h/2, f"{int(h):,}", ha="center", va="center")

def annotate_points(ax, x, y):
    for xi, yi in zip(x, y):
        if pd.notna(yi):
            ax.annotate(f"{yi:.2f}", (xi, yi), textcoords="offset points", xytext=(0,6), ha="center", fontsize=8)

# ===================== DATA SOURCE =====================
st.sidebar.header("Data source")
FILENAME = "OnTopSLA_Clean_from_orgnl.xlsx"

def find_raw_file(filename: str):
    here = Path(__file__).resolve()
    app_dir = here.parent      # dashboards/
    cwd = Path.cwd()

    raw_dirs = []
    raw_dirs += [app_dir / "RAW"]
    raw_dirs += [p / "RAW" for p in app_dir.parents[:4]]
    raw_dirs += [cwd / "RAW"]
    raw_dirs += [p / "RAW" for p in cwd.parents[:4]]

    seen = set()
    candidates = []
    for d in raw_dirs:
        p = (d / filename).resolve()
        if str(p) not in seen:
            candidates.append(p); seen.add(str(p))

    extra_bases = [app_dir, app_dir.parent, cwd]
    for base in extra_bases:
        if base.exists():
            try:
                for p in base.rglob(filename):
                    parts = [seg.lower() for seg in p.parts]
                    if "raw" in parts and str(p.resolve()) not in seen:
                        candidates.append(p.resolve()); seen.add(str(p.resolve()))
            except Exception:
                pass

    for c in candidates:
        if c.exists():
            return c, candidates
    return None, candidates

RAW_FOUND, RAW_CANDIDATES = find_raw_file(FILENAME)
uploaded = st.sidebar.file_uploader("Upload Excel (override RAW)", type=["xlsx"])

if uploaded is not None:
    data = load_clean_excel(uploaded.read())
elif RAW_FOUND is not None:
    st.sidebar.success(f"Usando RAW: {RAW_FOUND}")
    data = load_clean_excel(str(RAW_FOUND))
else:
    st.error("No se encontró el Excel en ninguna ruta esperada (RAW). Sube uno con el uploader.")
    with st.sidebar.expander("Rutas probadas", expanded=False):
        for c in RAW_CANDIDATES:
            st.write(str(c))
    st.stop()

# Normaliza semanas donde aplica
for key in ["OnTop_SLA_Global","OnTop_FRT_Global","SLA_By_Team","FRT_By_Team",
            "Tickets_Created","Tickets_Solved","Tickets_Breached","Tickets_Unbreached",
            "Tickets_Inflow","Tickets_Breach_Split","mrt_global","mrt_by_team"]:
    data[key] = enforce_int_week(data.get(key, pd.DataFrame()))

# Normaliza categorías
category_df = normalize_category_df(data.get("Category_By_Week"))

# Target SLA
target = infer_target(data.get("Summary"), SLA_TARGET_DEFAULT)

# Marts (si faltan, los armamos)
mrt_global = data.get("mrt_global", pd.DataFrame())
if mrt_global is None or mrt_global.empty:
    frt = data.get("OnTop_FRT_Global", pd.DataFrame()).copy()
    sla = data.get("OnTop_SLA_Global", pd.DataFrame()).copy()
    if not sla.empty: sla.columns = [c.lower() for c in sla.columns]
    if not frt.empty: frt.columns = [c.lower() for c in frt.columns]
    if {"year","week","sla_attainment_pct"}.issubset(sla.columns) and {"year","week","frt_business_hours"}.issubset(frt.columns):
        mrt_global = pd.merge(sla[["year","week","sla_attainment_pct"]],
                              frt[["year","week","frt_business_hours"]],
                              on=["year","week"], how="outer")
else:
    mrt_global.columns = [c.lower() for c in mrt_global.columns]

mrt_by_team = data.get("mrt_by_team", pd.DataFrame())
if mrt_by_team is None or mrt_by_team.empty:
    sla_t = data.get("SLA_By_Team", pd.DataFrame()).copy()
    frt_t = data.get("FRT_By_Team", pd.DataFrame()).copy()
    if not sla_t.empty: sla_t.columns = [c.lower() for c in sla_t.columns]
    if not frt_t.empty: frt_t.columns = [c.lower() for c in frt_t.columns]
    need_s = {"team","year","week","sla_attainment_pct"}.issubset(sla_t.columns)
    need_f = {"team","year","week","frt_business_hours"}.issubset(frt_t.columns)
    if need_s and need_f:
        mrt_by_team = pd.merge(sla_t[["team","year","week","sla_attainment_pct"]],
                               frt_t[["team","year","week","frt_business_hours"]],
                               on=["team","year","week"], how="outer")
else:
    mrt_by_team.columns = [c.lower() for c in mrt_by_team.columns]

# ===================== FILTERS =====================
st.sidebar.header("Filters")

def week_universe():
    cols = []
    for k in ["mrt_global","mrt_by_team","Tickets_Created","Tickets_Solved","Tickets_Breached","Tickets_Unbreached", "Category_By_Week"]:
        df = category_df if k=="Category_By_Week" else data.get(k, pd.DataFrame())
        if not isinstance(df, pd.DataFrame) or df.empty or "week" not in df:
            continue
        cols += pd.to_numeric(df["week"], errors="coerce").dropna().astype(int).tolist()
    return sorted(set(cols))

years = sorted(set(pd.concat([
    mrt_global.get("year", pd.Series(dtype=int)),
    mrt_by_team.get("year", pd.Series(dtype=int)),
    data.get("Tickets_Created", pd.DataFrame()).get("year", pd.Series(dtype=int)),
    data.get("Tickets_Solved", pd.DataFrame()).get("year", pd.Series(dtype=int)),
    category_df.get("year", pd.Series(dtype=int)),
], ignore_index=True).dropna().astype(int).unique().tolist())) or [2025]

year = st.sidebar.selectbox("Year", years, index=len(years)-1)
weeks_avail = week_universe() or list(range(1,53))
wk_min, wk_max = st.sidebar.select_slider("Weeks range", options=weeks_avail, value=(min(weeks_avail), max(weeks_avail)))
teams = sorted(mrt_by_team.get("team", pd.Series(dtype=str)).dropna().unique().tolist())
selected_teams = st.sidebar.multiselect("Teams", teams, default=teams)

# Filtrado de marts
mg = mrt_global.copy()
if not mg.empty:
    mg = mg[(mg["year"]==year) & (mg["week"]>=wk_min) & (mg["week"]<=wk_max)]
# --- copia sin filtrar para usarla en "Teams Under Target" (org-wide) ---
mrt_by_team_full = mrt_by_team.copy()
mbt = mrt_by_team.copy()
if not mbt.empty:
    mbt = mbt[(mbt["year"]==year) & (mbt["week"]>=wk_min) & (mbt["week"]<=wk_max)]
    if selected_teams:
        mbt = mbt[mbt["team"].isin(selected_teams)]

# Escala % si venía 0–1
if "sla_attainment_pct" in mg:
    mg["sla_attainment_pct"] = scale_pct_if_needed(mg["sla_attainment_pct"])
if "sla_attainment_pct" in mbt:
    mbt["sla_attainment_pct"] = scale_pct_if_needed(mbt["sla_attainment_pct"])

# ===================== VERSIONES FILTRADAS para inflows/breaches/categorías =====================
tc_full = data.get("Tickets_Created", pd.DataFrame())
ts_full = data.get("Tickets_Solved", pd.DataFrame())
tb_full = data.get("Tickets_Breached", pd.DataFrame())
tu_full = data.get("Tickets_Unbreached", pd.DataFrame())
cat_full = category_df if isinstance(category_df, pd.DataFrame) else pd.DataFrame()

def _filter_year_weeks(df):
    if not isinstance(df, pd.DataFrame) or df.empty:
        return pd.DataFrame()
    if "year" not in df or "week" not in df:
        return pd.DataFrame()
    out = df.copy()
    out["week"] = pd.to_numeric(out["week"], errors="coerce")
    out = out.dropna(subset=["week"])
    out["week"] = out["week"].astype(int)
    return out[(out["year"] == year) & (out["week"] >= wk_min) & (out["week"] <= wk_max)]

tc = _filter_year_weeks(tc_full)
ts = _filter_year_weeks(ts_full)
tb = _filter_year_weeks(tb_full)
tu = _filter_year_weeks(tu_full)
cat_f = _filter_year_weeks(cat_full)

# ===================== KPIs =====================
st.title("OnTop SLA Dashboard")
st.caption(f"Business Hours: 10h/day (08:00–18:00 America/Bogota). Target SLA: {target:.0f}%.")

c1,c2,c3,c4 = st.columns(4)

sla_global_avg = float(mg["sla_attainment_pct"].mean()) if ("sla_attainment_pct" in mg) and not mg.empty else np.nan
frt_global_avg = float(mg["frt_business_hours"].mean()) if ("frt_business_hours" in mg) and not mg.empty else np.nan
gap = (sla_global_avg - target) if not np.isnan(sla_global_avg) else np.nan

current_week = wk_max
sla_series = mg.set_index("week")["sla_attainment_pct"] if ("sla_attainment_pct" in mg) and not mg.empty else pd.Series(dtype=float)
frt_series = mg.set_index("week")["frt_business_hours"] if ("frt_business_hours" in mg) and not mg.empty else pd.Series(dtype=float)
sla_curr, sla_prev, sla_delta = wow_delta(sla_series, current_week) if not sla_series.empty else (np.nan,np.nan,np.nan)
frt_curr, frt_prev, frt_delta = wow_delta(frt_series, current_week) if not frt_series.empty else (np.nan,np.nan,np.nan)

# Team compliance (>= target) última semana del rango
team_compliance_curr = np.nan; team_compliance_prev = np.nan; team_compliance_delta = np.nan
if not mbt.empty and "sla_attainment_pct" in mbt:
    tmp = mbt.copy()
    tmp["within_target"] = np.where(tmp["sla_attainment_pct"] >= target, 1, 0)
    comp = tmp.groupby("week")["within_target"].agg(["sum","count"])
    if current_week in comp.index and comp.loc[current_week,"count"]>0:
        team_compliance_curr = comp.loc[current_week,"sum"]*100/comp.loc[current_week,"count"]
    if (current_week-1) in comp.index and comp.loc[current_week-1,"count"]>0:
        team_compliance_prev = comp.loc[current_week-1,"sum"]*100/comp.loc[current_week-1,"count"]
    if pd.notna(team_compliance_curr) and pd.notna(team_compliance_prev):
        team_compliance_delta = team_compliance_curr - team_compliance_prev

c1.metric("SLA Global (avg %)", f"{sla_global_avg:.2f}%" if not np.isnan(sla_global_avg) else "NA",
          delta=f"{gap:+.2f} pp" if not np.isnan(gap) else None)
c2.metric(f"SLA Week {current_week}", f"{sla_curr:.2f}%" if not np.isnan(sla_curr) else "NA",
          delta=f"{sla_delta:+.2f} pp" if not np.isnan(sla_delta) else None)
c3.metric(f"FRT Week {current_week} (hrs)", f"{frt_curr:.2f}" if not np.isnan(frt_curr) else "NA",
          delta=f"{frt_delta:+.2f}" if not np.isnan(frt_delta) else None)
c4.metric(f"Team Compliance Week {current_week}", f"{team_compliance_curr:.2f}%" if not np.isnan(team_compliance_curr) else "NA",
          delta=f"{team_compliance_delta:+.2f} pp" if not np.isnan(team_compliance_delta) else None)

st.markdown("---")

# ===================== AI NARRATIVE INSIGHTS (Local NLG, EN, icons, neutral endings) =====================
st.subheader("AI Narrative Insights (Local NLG)")

# --- Styles (bigger text + color cues) ---
st.markdown("""
<style>
.narrative p, .narrative li { font-size:1.08rem; line-height:1.6; color:#111; }
.narrative .up   { color:#1a9850; font-weight:700; }   /* green = positive */
.narrative .down { color:#d73027; font-weight:700; }   /* red = negative */
.narrative .muted{ color:#777; }
.narrative h4    { margin: 0.3rem 0 0.5rem 0; font-size:1.1rem; }
</style>
""", unsafe_allow_html=True)

# ---------- helpers ----------
def _pct_change(new, old):
    if pd.isna(new) or pd.isna(old) or old == 0:
        return np.nan
    return (new - old) * 100.0 / old

def _fmt_pp(x):
    return "NA" if pd.isna(x) else f"{x:+.2f} pp"

def _fmt_pct(x):
    return "NA" if pd.isna(x) else f"{x:+.1f}%"

def _fmt_hr(x):
    return "NA" if pd.isna(x) else f"{x:+.2f} hrs" if x < 0 or str(x).startswith(("+","-")) else f"{x:.2f} hrs"

# ---------- main generator ----------
def generate_narrative_insights_en(mg_f: pd.DataFrame,
                                   mbt_f: pd.DataFrame,
                                   target: float,
                                   wk_min: int,
                                   wk_max: int) -> str:
    """
    Builds executive narrative using only local dataframes (no API).
    Expects mg_f with columns: ['week','sla_attainment_pct','frt_business_hours']
            mbt_f with columns: ['team','week','sla_attainment_pct'] (optional for team breakdown)
    """
    if mg_f is None or mg_f.empty or "week" not in mg_f:
        return "<div class='narrative'><p><em>Not enough data to generate narrative insights.</em></p></div>"

    # --- Context: last week with data inside the selected window ---
    weeks_num = pd.to_numeric(mg_f["week"], errors="coerce").dropna().astype(int)
    if weeks_num.empty:
        return "<div class='narrative'><p><em>Not enough data to generate narrative insights.</em></p></div>"
    last_w = int(weeks_num.max())

    g = mg_f.set_index("week")

    sla_w   = g.get("sla_attainment_pct", pd.Series(dtype=float)).get(last_w,   np.nan)
    sla_wm1 = g.get("sla_attainment_pct", pd.Series(dtype=float)).get(last_w-1, np.nan)
    frt_w   = g.get("frt_business_hours", pd.Series(dtype=float)).get(last_w,   np.nan)
    frt_wm1 = g.get("frt_business_hours", pd.Series(dtype=float)).get(last_w-1, np.nan)

    d_sla    = sla_w - sla_wm1 if pd.notna(sla_w) and pd.notna(sla_wm1) else np.nan
    d_frt    = frt_w - frt_wm1 if pd.notna(frt_w) and pd.notna(frt_wm1) else np.nan
    d_frtpct = _pct_change(frt_w, frt_wm1)

    # performance colors
    cls_sla_lvl = "up" if (pd.notna(sla_w) and sla_w >= target) else "down"
    cls_frt_lvl = "up" if (pd.notna(d_frt) and d_frt < 0) else "down"   # lowering FRT is good

    # ---------------------------
    # What moved the needle (SLA)
    # ---------------------------
    moved_lines = []
    if pd.notna(d_sla) and abs(d_sla) >= 0.1:
        if pd.notna(d_frt) and ((d_sla > 0 and d_frt < 0) or (d_sla < 0 and d_frt > 0)):
            reason = "a reduction" if d_frt < 0 else "an increase"
            moved_lines.append(
                f"<li><b>What moved the needle:</b> SLA changed "
                f"<span class='{cls_sla_lvl}'>{_fmt_pp(d_sla)}</span>, driven by {reason} in FRT of "
                f"<span class='{cls_frt_lvl}'>{_fmt_hr(d_frt)}</span> "
                f"({_fmt_pct(d_frtpct)} vs W{last_w-1}).</li>"
            )
        else:
            # fallback: teams with biggest WoW improvement
            top_up = []
            if (mbt_f is not None and not mbt_f.empty and
                {"team","week","sla_attainment_pct"}.issubset(mbt_f.columns)):
                cur  = (mbt_f[mbt_f["week"]==last_w][["team","sla_attainment_pct"]]
                        .rename(columns={"sla_attainment_pct":"curr"}))
                prev = (mbt_f[mbt_f["week"]==last_w-1][["team","sla_attainment_pct"]]
                        .rename(columns={"sla_attainment_pct":"prev"}))
                chg = cur.merge(prev, on="team", how="left")
                chg["delta"] = chg["curr"] - chg["prev"]
                chg = chg.dropna(subset=["delta"]).sort_values("delta", ascending=False).head(3)
                top_up = [f"{r.team} (+{r.delta:.1f} pp)" for _, r in chg.iterrows()]
            if top_up:
                moved_lines.append(
                    f"<li><b>What moved the needle:</b> SLA improvement mainly from teams: "
                    f"{', '.join(top_up)}.</li>"
                )
    if not moved_lines:
        moved_lines.append("<li><b>What moved the needle:</b> Minor week-over-week variation; no dominant driver identified.</li>")

    # ---------------------------------
    # Negative trends & risk areas
    # ---------------------------------
    risk_lines = []
    if pd.notna(sla_w) and sla_w < target:
        risk_lines.append(
            f"<li><b>SLA Risk:</b> Week {last_w} below goal "
            f"(<span class='down'>{sla_w:.2f}%</span> vs target {target:.0f}%).</li>"
        )
    if pd.notna(d_frt) and d_frt > 0.25:
        risk_lines.append(
            f"<li><b>FRT Increasing:</b> <span class='down'>{_fmt_hr(d_frt)}</span> vs W{last_w-1} "
            f"({_fmt_pct(d_frtpct)}).</li>"
        )

    # Teams below goal with conditional message
    if (mbt_f is not None and not mbt_f.empty and
        {"team","week","sla_attainment_pct"}.issubset(mbt_f.columns)):
        teams_last = (mbt_f[mbt_f["week"]==last_w][["team","sla_attainment_pct"]]
                      .dropna(subset=["sla_attainment_pct"]))
        under = teams_last[teams_last["sla_attainment_pct"] < target].sort_values("sla_attainment_pct")
        if not under.empty:
            n_under = int(len(under))
            n_total = int(teams_last["team"].nunique())
            worst = ", ".join([f"{r.team} ({r.sla_attainment_pct:.1f}%)" for _, r in under.head(3).iterrows()])

            # Message policy (neutral text color)
            if pd.notna(sla_w) and sla_w < target:
                tail = "Performance gaps in these teams are limiting SLA target achievement."
            else:
                if n_under <= 3:
                    tail = "The current performance may impact overall SLA compliance."
                elif n_under <= 5:
                    tail = "These results could lower overall SLA achievement if not corrected."
                else:
                    tail = "Performance gaps in these teams are limiting the SLA target achievement."

            risk_lines.append(
                f"<li><b>Teams below goal:</b> {n_under}/{n_total} "
                f"(lowest performers: {worst}). {tail}</li>"
            )

    if not risk_lines:
        risk_lines.append("<li><b>Risk Summary:</b> No significant risks this week.</li>")

    # ---------------------------------
    # Executive summary
    # ---------------------------------
    summary_line = (
        f"<p><b>Executive summary (W{last_w}):</b> "
        f"SLA: <span class='{cls_sla_lvl}'>{sla_w:.2f}%</span> "
        f"({_fmt_pp(d_sla)} vs W{last_w-1}); "
        f"FRT: <span class='{cls_frt_lvl}'>{_fmt_hr(frt_w)}</span> "
        f"({_fmt_hr(d_frt)} vs W{last_w-1}). "
        f"Target remains at <b>{target:.0f}%</b>.</p>"
    )

    html = f"""
    <div class='narrative'>
      <h4>📅 Analysis Window</h4>
      <p>Year {year}, weeks {wk_min}–{wk_max}. Latest week with data: <b>{last_w}</b>.</p>

      <h4>📈 What moved the needle</h4>
      <ul>
        {''.join(moved_lines)}
      </ul>

      <h4>⚠️ Negative trends & risk areas</h4>
      <ul>
        {''.join(risk_lines)}
      </ul>

      <h4>🧭 Executive summary</h4>
      {summary_line}
    </div>
    """
    return html

# ---------- render ----------
narr_html = generate_narrative_insights_en(mg, mbt, target, wk_min, wk_max)
st.markdown(narr_html, unsafe_allow_html=True)
st.markdown("---")

# ===================== CHARTS (compactos, 2 por fila) =====================
colA, colB = st.columns(2)

with colA:
    st.subheader("Global SLA Trend (Weekly)")
    if not mg.empty and "sla_attainment_pct" in mg:
        g = mg.dropna(subset=["week","sla_attainment_pct"]).sort_values("week").copy()
        g["week"] = g["week"].astype(int)
        fig, ax = plt.subplots(); fig.set_size_inches(6,3)
        ax.plot(g["week"], g["sla_attainment_pct"], marker="o")
        ax.axhline(y=target, linestyle="--")
        ax.set_xlabel("Week"); ax.set_ylabel("SLA %"); ax.set_xticks(sorted(g["week"].unique()))
        annotate_points(ax, g["week"].values, g["sla_attainment_pct"].values)
        st.pyplot(fig)
    else:
        st.info("No global SLA data for filters.")

with colB:
    st.subheader("Global FRT Trend (Weekly)")
    if not mg.empty and "frt_business_hours" in mg:
        g = mg.dropna(subset=["week","frt_business_hours"]).sort_values("week").copy()
        g["week"] = g["week"].astype(int)
        fig, ax = plt.subplots(); fig.set_size_inches(6,3)
        ax.plot(g["week"], g["frt_business_hours"], marker="o")
        ax.set_xlabel("Week"); ax.set_ylabel("FRT (Business Hours)"); ax.set_xticks(sorted(g["week"].unique()))
        annotate_points(ax, g["week"].values, g["frt_business_hours"].values)
        st.pyplot(fig)
    else:
        st.info("No global FRT data for filters.")

colC, colD = st.columns(2)

with colC:
    st.subheader("Inflows: Created vs Solved")
    if not tc.empty or not ts.empty:
        df = pd.merge(
            tc.rename(columns={"tickets_created":"created"}) if "tickets_created" in tc else tc,
            ts.rename(columns={"tickets_solved":"solved"}) if "tickets_solved" in ts else ts,
            on=["year","week"], how="outer"
        ).fillna(0)
        if not df.empty:
            df = df.sort_values("week")
            fig, ax = plt.subplots(); fig.set_size_inches(6,3)
            x = df["week"].astype(int).values
            created = pd.to_numeric(df.get("created", pd.Series([0]*len(df))), errors="coerce").fillna(0).values
            solved  = pd.to_numeric(df.get("solved",  pd.Series([0]*len(df))), errors="coerce").fillna(0).values
            ax.plot(x, created, marker="o", label="Created")
            ax.plot(x, solved,  marker="o", label="Solved")
            for xi, yi in zip(x, created):
                ax.annotate(f"{int(yi)}", (xi, yi), textcoords="offset points", xytext=(0,6), ha="center", fontsize=8)
            for xi, yi in zip(x, solved):
                ax.annotate(f"{int(yi)}", (xi, yi), textcoords="offset points", xytext=(0,6), ha="center", fontsize=8)
            ax.set_xlabel("Week"); ax.set_ylabel("Tickets"); ax.set_xticks(sorted(df["week"].astype(int).unique()))
            ax.legend()
            st.pyplot(fig)
        else:
            st.info("No inflow data for selected filters.")
    else:
        st.info("Tickets Created/Solved not available for selected filters.")

with colD:
    st.subheader("Breaches: Breached vs Unbreached")
    if not tb.empty or not tu.empty:
        df = pd.merge(tb, tu, on=["year","week"], how="outer").fillna(0)
        if not df.empty:
            df = df.sort_values("week")
            fig, ax = plt.subplots(); fig.set_size_inches(6,3)
            x = df["week"].astype(int).values
            bcol = next((c for c in df.columns if "breach" in c and "un" not in c), "tickets_breached")
            ucol = next((c for c in df.columns if "unbreach" in c or "on-time" in c or "ok" in c), "tickets_unbreached")
            b = pd.to_numeric(df.get(bcol, 0), errors="coerce").fillna(0).astype(int).values
            u = pd.to_numeric(df.get(ucol, 0), errors="coerce").fillna(0).astype(int).values
            ax.bar(x, b, label="Breached"); ax.bar(x, u, bottom=b, label="Unbreached")
            ax.set_xlabel("Week"); ax.set_ylabel("Tickets"); ax.set_xticks(sorted(df["week"].astype(int).unique()))
            annotate_bars(ax, x, b); annotate_stacked_bars(ax, x, b, u)
            ax.legend()
            st.pyplot(fig)
        else:
            st.info("No breach data for selected filters.")
    else:
        st.info("Tickets Breached/Unbreached not available.")

# ===================== CATEGORY STACKED (Top-5 última semana del rango) =====================
st.markdown("---")
st.subheader("Contact Categories by Week (Top-5 of last week, Stacked)")

if not cat_f.empty:
    csel = cat_f.copy()
    if not csel.empty and {"category","tickets"}.issubset(csel.columns):
        weeks_present = csel["week"].dropna().astype(int)
        lastw = int(weeks_present.max())
        top5 = (csel[csel["week"]==lastw].groupby("category", as_index=False)["tickets"].sum()
                    .sort_values("tickets", ascending=False).head(5)["category"].tolist())
        ctop = csel[csel["category"].isin(top5)].copy()
        pivot = ctop.pivot_table(index="week", columns="category", values="tickets", aggfunc="sum").fillna(0)
        pivot = pivot.sort_index()
        fig, ax = plt.subplots(); fig.set_size_inches(12, 4)
        x = pivot.index.astype(int).values
        bottom = np.zeros(len(pivot))
        for col in pivot.columns:
            y = pivot[col].values
            ax.bar(x, y, bottom=bottom, label=col)
            for i, (xx, yy, bb) in enumerate(zip(x, y, bottom)):
                if yy > 0:
                    ax.text(xx, bb + yy/2, f"{int(yy)}", ha="center", va="center", fontsize=8)
            bottom += y
        ax.set_xlabel("Week"); ax.set_ylabel("Tickets")
        ax.set_xticks(x)
        ax.legend(ncol=3, fontsize=8, bbox_to_anchor=(1.0, 1.02), loc="lower right")
        st.caption(f"Top-5 calculado a partir de la semana {lastw}")
        st.pyplot(fig)
    else:
        st.info("No category data for the selected weeks.")
else:
    st.info("Category sheet not available or empty.")

st.markdown("---")

# ===================== HEATMAP (Tramos) & TEAMS UNDER TARGET (última semana del rango) =====================
st.subheader("SLA by Team and Week (Heatmap + values)")
if not mbt.empty and {"team","week","sla_attainment_pct"}.issubset(mbt.columns):
    t = mbt.dropna(subset=["team","week","sla_attainment_pct"]).copy()
    t["week"] = t["week"].astype(int)
    t["sla_attainment_pct"] = pd.to_numeric(t["sla_attainment_pct"], errors="coerce").clip(lower=0, upper=100)

    pivot = t.pivot_table(index="team", columns="week", values="sla_attainment_pct", aggfunc="mean")
    pivot = pivot.sort_index().sort_index(axis=1)
    if pivot.size > 0:
        fig, ax = plt.subplots(); fig.set_size_inches(12, 6)
        boundaries = [0, 60, 80, 95, 100]
        colors = ["#d73027", "#fc8d59", "#fee08b", "#1a9850"]  # rojo, naranja, amarillo, verde
        cmap = ListedColormap(colors)
        norm = BoundaryNorm(boundaries, cmap.N, clip=True)
        im = ax.imshow(pivot.values.astype(float), aspect="auto", cmap=cmap, norm=norm)
        ax.set_yticks(range(len(pivot.index))); ax.set_yticklabels(pivot.index)
        ax.set_xticks(range(len(pivot.columns))); ax.set_xticklabels(pivot.columns)
        cbar = plt.colorbar(im, ax=ax, ticks=[30,70,87,97])
        cbar.set_label("SLA % (tramos: rojo<60, naranja<80, amarillo<95, verde≥95)")
        for i in range(len(pivot.index)):
            for j in range(len(pivot.columns)):
                val = pivot.values[i, j]
                if np.isnan(val): 
                    continue
                text_color = "black" if val >= 80 else "white"
                ax.text(j, i, f"{val:.1f}", ha="center", va="center", color=text_color, fontsize=9)
        st.pyplot(fig)
    else:
        st.info("No team-level SLA data after filtering.")
else:
    st.info("No team-level SLA data for selected filters.")

# ===================== EXPORTS (PPTX / PDF) =====================
st.markdown("---")
st.subheader("Exports")

EXPORT_DIR = Path("./exports")
EXPORT_DIR.mkdir(exist_ok=True)

def fig_sla(mg, target):
    fig, ax = plt.subplots(); fig.set_size_inches(10,4)
    if not mg.empty and "sla_attainment_pct" in mg:
        g = mg.dropna(subset=["week","sla_attainment_pct"]).sort_values("week")
        ax.plot(g["week"].astype(int), g["sla_attainment_pct"], marker="o")
        ax.axhline(y=target, linestyle="--")
        ax.set_xlabel("Week"); ax.set_ylabel("SLA %"); ax.set_title("Global SLA Trend (Weekly)")
        annotate_points(ax, g["week"].astype(int).values, g["sla_attainment_pct"].values)
    return fig

def fig_frt(mg):
    fig, ax = plt.subplots(); fig.set_size_inches(10,4)
    if not mg.empty and "frt_business_hours" in mg:
        g = mg.dropna(subset=["week","frt_business_hours"]).sort_values("week")
        ax.plot(g["week"].astype(int), g["frt_business_hours"], marker="o")
        ax.set_xlabel("Week"); ax.set_ylabel("FRT (Business Hours)"); ax.set_title("Global FRT Trend (Weekly)")
        annotate_points(ax, g["week"].astype(int).values, g["frt_business_hours"].values)
    return fig

def fig_inflow(tc, ts):
    fig, ax = plt.subplots(); fig.set_size_inches(10,4)
    if not tc.empty or not ts.empty:
        df = pd.merge(
            tc.rename(columns={"tickets_created":"created"}) if "tickets_created" in tc else tc,
            ts.rename(columns={"tickets_solved":"solved"}) if "tickets_solved" in ts else ts,
            on=["year","week"], how="outer"
        ).fillna(0).sort_values("week")
        x = df["week"].astype(int).values
        created = pd.to_numeric(df.get("created", 0), errors="coerce").fillna(0).values
        solved  = pd.to_numeric(df.get("solved",  0), errors="coerce").fillna(0).values
        ax.plot(x, created, marker="o", label="Created")
        ax.plot(x, solved,  marker="o", label="Solved")
        for xi, yi in zip(x, created):
            if pd.notna(yi): ax.annotate(f"{int(yi)}", (xi, yi), textcoords="offset points", xytext=(0,6), ha="center", fontsize=8)
        for xi, yi in zip(x, solved):
            if pd.notna(yi): ax.annotate(f"{int(yi)}", (xi, yi), textcoords="offset points", xytext=(0,6), ha="center", fontsize=8)
        ax.set_xlabel("Week"); ax.set_ylabel("Tickets"); ax.set_title("Inflows: Created vs Solved"); ax.legend()
    return fig

def fig_breaches(tb, tu):
    fig, ax = plt.subplots(); fig.set_size_inches(10,4)
    if not tb.empty or not tu.empty:
        df = pd.merge(tb, tu, on=["year","week"], how="outer").fillna(0).sort_values("week")
        x = df["week"].astype(int).values
        bcol = next((c for c in df.columns if "breach" in c and "un" not in c), "tickets_breached")
        ucol = next((c for c in df.columns if "unbreach" in c or "on-time" in c or "ok" in c), "tickets_unbreached")
        b = pd.to_numeric(df.get(bcol, 0), errors="coerce").fillna(0).astype(int).values
        u = pd.to_numeric(df.get(ucol, 0), errors="coerce").fillna(0).astype(int).values
        ax.bar(x, b, label="Breached"); ax.bar(x, u, bottom=b, label="Unbreached")
        for i, X in enumerate(x):
            if b[i] > 0: ax.text(X, b[i]/2, str(b[i]), ha="center", va="center", fontsize=8)
            if u[i] > 0: ax.text(X, b[i]+u[i]/2, str(u[i]), ha="center", va="center", fontsize=8)
        ax.set_xlabel("Week"); ax.set_ylabel("Tickets"); ax.set_title("Breached vs Unbreached"); ax.legend()
    return fig

def fig_category_top5(cat):
    fig, ax = plt.subplots(); fig.set_size_inches(10,4)
    if isinstance(cat, pd.DataFrame) and not cat.empty:
        csel = cat.copy()
        weeks_present = csel["week"].dropna().astype(int)
        if not weeks_present.empty and {"category","tickets"}.issubset(csel.columns):
            lastw = int(weeks_present.max())
            top5 = (csel[csel["week"]==lastw].groupby("category", as_index=False)["tickets"].sum()
                        .sort_values("tickets", ascending=False).head(5)["category"].tolist())
            ctop = csel[csel["category"].isin(top5)]
            pivot = ctop.pivot_table(index="week", columns="category", values="tickets", aggfunc="sum").fillna(0)
            pivot = pivot.sort_index(); x = pivot.index.astype(int).values
            bottom = np.zeros(len(pivot))
            for col in pivot.columns:
                y = pivot[col].values
                ax.bar(x, y, bottom=bottom, label=col)
                for i, (xx, yy, bb) in enumerate(zip(x, y, bottom)):
                    if yy > 0:
                        ax.text(xx, bb + yy/2, f"{int(yy)}", ha="center", va="center", fontsize=8)
                bottom += y
            ax.set_xlabel("Week"); ax.set_ylabel("Tickets"); ax.set_title(f"Categories by Week (Top 5 of week {lastw})")
            ax.legend(ncol=3, fontsize=8, bbox_to_anchor=(1.0, 1.02), loc="lower right")
    return fig

def fig_heatmap_discrete(mbt, target):
    fig, ax = plt.subplots(); fig.set_size_inches(12,6)
    if not mbt.empty and "sla_attainment_pct" in mbt:
        t = mbt.dropna(subset=["team","week","sla_attainment_pct"]).copy()
        t["week"] = t["week"].astype(int)
        t["sla_attainment_pct"] = pd.to_numeric(t["sla_attainment_pct"], errors="coerce").clip(lower=0, upper=100)
        pivot = t.pivot_table(index="team", columns="week", values="sla_attainment_pct", aggfunc="mean")
        pivot = pivot.sort_index().sort_index(axis=1)
        if pivot.size>0:
            boundaries = [0, 60, 80, 95, 100]
            colors = ["#d73027", "#fc8d59", "#fee08b", "#1a9850"]
            cmap = ListedColormap(colors)
            norm = BoundaryNorm(boundaries, cmap.N, clip=True)
            im = ax.imshow(pivot.values.astype(float), aspect="auto", cmap=cmap, norm=norm)
            ax.set_yticks(range(len(pivot.index))); ax.set_yticklabels(pivot.index)
            ax.set_xticks(range(len(pivot.columns))); ax.set_xticklabels(pivot.columns)
            plt.colorbar(im, ax=ax, ticks=[30,70,87,97], label="SLA % (tramos)")
            for i in range(len(pivot.index)):
                for j in range(len(pivot.columns)):
                    val = pivot.values[i, j]
                    if np.isnan(val): 
                        continue
                    color = "black" if val >= 80 else "white"
                    ax.text(j, i, f"{val:.1f}", ha="center", va="center", color=color, fontsize=9)
            ax.set_title("SLA by Team and Week")
    return fig

# Botones de export
col1, col2 = st.columns(2)
with col1:
    if st.button("Export PPTX"):
        try:
            base = f"{year}_w{wk_min}-{wk_max}"
            # Guardamos imágenes temporales
            figs = [
                ("sla.png",        fig_sla(mg, target)),
                ("frt.png",        fig_frt(mg)),
                ("inflow.png",     fig_inflow(tc, ts)),
                ("breaches.png",   fig_breaches(tb, tu)),
                ("category.png",   fig_category_top5(cat_f)),
                ("heatmap.png",    fig_heatmap_discrete(mbt, target)),
            ]
            paths = []
            for fname, fig in figs:
                out = EXPORT_DIR / f"{base}_{fname}"
                fig.savefig(out, bbox_inches="tight", dpi=150)
                plt.close(fig)
                paths.append(out)

            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "OnTop SLA Weekly Report"
            title_slide.placeholders[1].text = f"Year {year} | Weeks {wk_min}–{wk_max} | Target {target:.0f}%"
            for img in paths:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                left = Inches(0.5); top = Inches(0.8); width = Inches(9)
                slide.shapes.add_picture(str(img), left, top, width=width)
            pptx_path = EXPORT_DIR / f"OnTop_SLA_{base}.pptx"
            prs.save(pptx_path)
            st.success(f"PPTX generado: {pptx_path}")
            st.download_button("Descargar PPTX", data=open(pptx_path, "rb"), file_name=pptx_path.name)
        except Exception as e:
            st.error(f"Error generando PPTX: {e}")

with col2:
    if st.button("Export PDF"):
        try:
            base = f"{year}_w{wk_min}-{wk_max}"
            figs = [
                ("sla.png",        fig_sla(mg, target)),
                ("frt.png",        fig_frt(mg)),
                ("inflow.png",     fig_inflow(tc, ts)),
                ("breaches.png",   fig_breaches(tb, tu)),
                ("category.png",   fig_category_top5(cat_f)),
                ("heatmap.png",    fig_heatmap_discrete(mbt, target)),
            ]
            paths = []
            for fname, fig in figs:
                out = EXPORT_DIR / f"{base}_{fname}"
                fig.savefig(out, bbox_inches="tight", dpi=150)
                plt.close(fig)
                paths.append(out)

            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            from reportlab.lib.utils import ImageReader
            pdf_path = EXPORT_DIR / f"OnTop_SLA_{base}.pdf"
            c = canvas.Canvas(str(pdf_path), pagesize=letter)
            W, H = letter
            c.setFont("Helvetica-Bold", 18); c.drawString(50, H-80, "OnTop SLA Weekly Report")
            c.setFont("Helvetica", 12); c.drawString(50, H-100, f"Year {year} | Weeks {wk_min}-{wk_max} | Target {target:.0f}%")
            c.showPage()
            for img in paths:
                ir = ImageReader(str(img))
                c.drawImage(ir, 30, 80, width=W-60, height=H-160, preserveAspectRatio=True, anchor='c')
                c.showPage()
            c.save()
            st.success(f"PDF generado: {pdf_path}")
            st.download_button("Descargar PDF", data=open(pdf_path, "rb"), file_name=pdf_path.name)
        except Exception as e:
            st.error(f"Error generando PDF: {e}")